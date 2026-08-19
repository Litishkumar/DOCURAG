from dotenv import load_dotenv
import pathlib
_env_path = pathlib.Path(__file__).resolve().parent.parent.parent / '.env'
load_dotenv(dotenv_path=str(_env_path) if _env_path.exists() else None)
from flask import Flask, request, jsonify, send_from_directory, redirect
from flask_cors import CORS

from langchain.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.pgvector import PGVector
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import google.generativeai as genai
import os, uuid, warnings, shutil, re, json

from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash

from pptx import Presentation
import base64
import logging
import seaborn as sns
import requests
from bs4 import BeautifulSoup
from datetime import datetime, timedelta, timezone

import sqlalchemy
from sqlalchemy import create_engine, text

import jwt as pyjwt

# Configure matplotlib for headless environment
plt.switch_backend('Agg')
plt.style.use('dark_background')

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_AI_API_KEY")
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    logger.info("✓ Gemini API configured")
else:
    warnings.warn("⚠️ GEMINI_API_KEY not set")

# Postgres configuration (for chat history + pgvector document storage)
DATABASE_URL = os.getenv("DATABASE_URL")  # e.g. postgresql://user:pass@host:5432/dbname
PG_CONNECTION_STRING = None
pg_engine = None

if DATABASE_URL:
    try:
        # SQLAlchemy engine for our own tables (chat sessions/messages)
        pg_engine = create_engine(DATABASE_URL)
        with pg_engine.connect() as conn:
            conn.execute(text("""
                CREATE TABLE IF NOT EXISTS chat_sessions (
                    id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                    title TEXT NOT NULL DEFAULT 'New Chat',
                    created_at TIMESTAMP DEFAULT NOW()
                );
            """))
            conn.execute(text("""
                CREATE TABLE IF NOT EXISTS chat_messages (
                    id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                    session_id UUID NOT NULL REFERENCES chat_sessions(id) ON DELETE CASCADE,
                    role TEXT NOT NULL,
                    content TEXT NOT NULL,
                    mode TEXT,
                    sources TEXT,
                    chart_url TEXT,
                    created_at TIMESTAMP DEFAULT NOW()
                );
            """))
            # gen_random_uuid() needs pgcrypto on some Postgres setups; Supabase has it enabled by default.
            conn.commit()
        # PGVector uses its own SQLAlchemy-style connection string (psycopg2 driver)
        PG_CONNECTION_STRING = DATABASE_URL.replace("postgresql://", "postgresql+psycopg2://")
        logger.info("✓ Postgres connected — chat_sessions & chat_messages ready")
    except Exception as e:
        logger.error(f"⚠️ Postgres connection failed, falling back to local-only mode: {str(e)}")
        pg_engine = None
else:
    logger.warning("⚠️ DATABASE_URL not set — chat history & pgvector disabled, using local ChromaDB only")

# JWT Configuration
JWT_SECRET = os.getenv("JWT_SECRET", "docurag-jwt-secret-" + (DATABASE_URL or "local")[:16])

# Google OAuth Configuration (optional)
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET")
GOOGLE_REDIRECT_URI = os.getenv("GOOGLE_REDIRECT_URI", "http://localhost:5000/auth/google/callback")

# Create users table
if pg_engine:
    try:
        with pg_engine.connect() as conn:
            conn.execute(text("""
                CREATE TABLE IF NOT EXISTS users (
                    id UUID PRIMARY KEY DEFAULT gen_random_uuid(),
                    fullname TEXT,
                    username TEXT UNIQUE,
                    email TEXT UNIQUE NOT NULL,
                    password_hash TEXT,
                    picture TEXT,
                    provider TEXT DEFAULT 'email',
                    created_at TIMESTAMP DEFAULT NOW()
                );
            """))
            conn.commit()
        logger.info("✓ Users table ready")
    except Exception as e:
        logger.error(f"⚠️ Users table creation failed: {str(e)}")

# Document Loaders
LOADERS = {}
try:
    from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
    LOADERS.update({'pdf': PyPDFLoader, 'txt': TextLoader, 'docx': Docx2txtLoader})
    logger.info("✓ Document loaders loaded successfully")
except ImportError as e:
    logger.warning("⚠️ Limited loader support")

# Flask setup
app = Flask(__name__)
CORS(app, resources={r"/*": {
    "origins": [
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "http://localhost:5000",
        "http://127.0.0.1:5000",
    ],
    "allow_headers": ["Authorization", "Content-Type", "Accept"],
    "methods": ["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    "supports_credentials": True
}})

# Directories
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'uploads')
PLOTS_FOLDER = os.path.join(os.path.dirname(__file__), 'static', 'plots')
VECTOR_DB_DIR = os.path.join(os.path.dirname(__file__), 'db_miniLM')

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PLOTS_FOLDER, exist_ok=True)
os.makedirs(os.path.dirname(PLOTS_FOLDER), exist_ok=True)  # Ensure static folder exists

# Constants
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'pptx', 'ppt', 'xlsx', 'xls', 'csv'}
excel_data_store = {}
_data_summary_cache = {}  # {filename: summary_string} — avoids recomputing on each question
last_uploaded_type = None  # tracks 'data' or 'document' — whichever was uploaded most recently
_documents_indexed = False  # fast flag to avoid expensive similarity_search checks

# Keywords
GRAPH_KEYWORDS = ['graph', 'chart', 'plot', 'visualize', 'visualization', 'bar', 'line', 'pie', 'scatter', 'histogram', 'show', 'draw', 'create chart']
DOCUMENT_KEYWORDS = ['summarise', 'summarize', 'summary', 'document', 'content', 'explain', 'analyze', 'findings']

# Helper Functions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def using_postgres():
    return pg_engine is not None and PG_CONNECTION_STRING is not None

_embeddings_cache = None

def get_embeddings():
    """Loads the embedding model once and reuses it."""
    global _embeddings_cache
    if _embeddings_cache is None:
        logger.info("⏳ Loading embedding model (cached after this)...")
        _embeddings_cache = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
    return _embeddings_cache

# Pre-warm embedding model at import time so first request isn't slow
#try:
    #get_embeddings()
#except Exception as _e:
   # logger.warning(f"⚠️ Embedding model pre-warm failed (will retry on first use): {_e}")

def get_vector_store():
    """Returns a PGVector store if Postgres is configured, otherwise falls back
    to local ChromaDB. This is the single source of truth for document storage."""
    if using_postgres():
        return PGVector(
            connection_string=PG_CONNECTION_STRING,
            embedding_function=get_embeddings(),
            collection_name="documents"
        )
    return Chroma(persist_directory=VECTOR_DB_DIR, embedding_function=get_embeddings())

def index_documents(texts):
    """Adds document chunks to whichever vector store is active."""
    if using_postgres():
        PGVector.from_documents(
            documents=texts,
            embedding=get_embeddings(),
            connection_string=PG_CONNECTION_STRING,
            collection_name="documents"
        )
    else:
        db = Chroma.from_documents(texts, get_embeddings(), persist_directory=VECTOR_DB_DIR)
        db.persist()

def sparse_search(question, k=5):
    """Keyword-based search using Postgres full-text search (ts_rank) over the
    same document chunks PGVector stores. This is the 'sparse' half of hybrid
    retrieval — it catches exact terms, names, and numbers that dense/semantic
    search can sometimes blur past. Returns [] if not using Postgres."""
    if not using_postgres():
        return []
    try:
        with pg_engine.connect() as conn:
            rows = conn.execute(text("""
                SELECT e.document, e.cmetadata,
                       ts_rank(to_tsvector('english', e.document), plainto_tsquery('english', :q)) AS rank
                FROM langchain_pg_embedding e
                JOIN langchain_pg_collection c ON e.collection_id = c.uuid
                WHERE c.name = 'documents'
                  AND to_tsvector('english', e.document) @@ plainto_tsquery('english', :q)
                ORDER BY rank DESC
                LIMIT :k
            """), {'q': question, 'k': k}).fetchall()
        return [(row[0], row[1]) for row in rows]  # (document_text, metadata)
    except Exception as e:
        logger.error(f"⚠️ Sparse search failed: {str(e)}")
        return []

def hybrid_search(question, k=5):
    """Combines dense (semantic/PGVector) and sparse (keyword/Postgres full-text)
    search results using Reciprocal Rank Fusion, then returns the top k merged
    results. Falls back to dense-only if Postgres/sparse search isn't available."""
    dense_docs = get_vector_store().similarity_search(question, k=k)
    sparse_results = sparse_search(question, k=k)

    if not sparse_results:
        # No sparse results (no Postgres, or no keyword matches) — dense-only is fine
        return dense_docs

    # Reciprocal Rank Fusion: score = sum of 1/(60 + rank) across both lists
    RRF_K = 60
    scores = {}
    doc_lookup = {}

    for rank, doc in enumerate(dense_docs):
        key = doc.page_content[:200]  # dedupe key
        scores[key] = scores.get(key, 0) + 1.0 / (RRF_K + rank + 1)
        doc_lookup[key] = doc

    from langchain_core.documents import Document as LC_Document
    for rank, (doc_text, metadata) in enumerate(sparse_results):
        key = doc_text[:200]
        scores[key] = scores.get(key, 0) + 1.0 / (RRF_K + rank + 1)
        if key not in doc_lookup:
            doc_lookup[key] = LC_Document(page_content=doc_text, metadata=metadata or {})

    ranked_keys = sorted(scores.keys(), key=lambda k: scores[k], reverse=True)
    return [doc_lookup[key] for key in ranked_keys[:k]]

def clear_vector_database():
    """Clears all indexed documents. With Postgres this is a clean DB delete
    (no file-lock issues). With Chroma this deletes the local folder."""
    global _documents_indexed
    _documents_indexed = False
    if using_postgres():
        try:
            with pg_engine.connect() as conn:
                conn.execute(text("DELETE FROM langchain_pg_embedding WHERE collection_id IN (SELECT uuid FROM langchain_pg_collection WHERE name = 'documents');"))
                conn.commit()
            logger.info("🗑️ Cleared Postgres document collection")
            return True
        except Exception as e:
            logger.error(f"⚠️ Could not clear Postgres vector data: {str(e)}")
            return False
    try:
        if os.path.exists(VECTOR_DB_DIR):
            shutil.rmtree(VECTOR_DB_DIR)
            logger.info("🗑️ Cleared local ChromaDB")
            return True
    except Exception as e:
        logger.error(f"⚠️ Could not clear vector database: {str(e)}")
        return False
    return True

def has_documents():
    global _documents_indexed
    if _documents_indexed:
        return True
    try:
        result = len(get_vector_store().similarity_search("test", k=1)) > 0
        if result:
            _documents_indexed = True
        return result
    except Exception as e:
        logger.error(f"⚠️ has_documents() check failed: {str(e)}")
        return False

def has_data():
    return len(excel_data_store) > 0

def is_graph_request(question):
    return any(keyword in question.lower() for keyword in GRAPH_KEYWORDS)

def is_document_specific_question(question):
    return any(keyword in question.lower() for keyword in DOCUMENT_KEYWORDS)

def search_web(query, max_results=3):
    """Simple web search using DuckDuckGo (you can replace with your preferred search API)"""
    try:
        # This is a simple implementation - replace with your preferred search API
        search_url = f"https://duckduckgo.com/html/?q={query}"
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(search_url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        results = []
        result_links = soup.find_all('a', class_='result__a')[:max_results]
        
        for link in result_links:
            title = link.get_text()
            url = link.get('href')
            if url and title:
                results.append({'title': title, 'url': url})
        
        return results
    except Exception as e:
        logger.error(f"Web search error: {str(e)}")
        return []

import time as _time

GEMINI_BACKUP_API_KEY = os.getenv("GEMINI_BACKUP_API_KEY")

def _is_rate_limit_error(error_str):
    """Check if an error is a rate-limit/quota issue (retryable with backoff)."""
    lower = error_str.lower()
    return any(marker in lower for marker in
               ["429", "quota", "rate limit", "rate_limit", "resource_exhausted",
                "too many requests", "rpm", "requests per minute"])

def _is_overload_error(error_str):
    """Check if an error is a transient overload (retryable after brief wait)."""
    lower = error_str.lower()
    return any(marker in lower for marker in
               ["503", "overloaded", "high demand", "unavailable", "timeout"])

def _is_model_not_found(error_str):
    """Check if the model name is invalid (skip immediately)."""
    lower = error_str.lower()
    return any(marker in lower for marker in
               ["not found", "404", "invalid model", "not supported", "does not exist"])

def ask_gemini(prompt, temp=0.7, include_web_context=False, web_query=None):
    if not GEMINI_API_KEY and not GEMINI_BACKUP_API_KEY:
        return "⚠️ Gemini API key not set"

    if include_web_context and web_query:
        web_results = search_web(web_query)
        if web_results:
            web_context = "Here are some relevant web results:\n"
            for i, r in enumerate(web_results, 1):
                web_context += f"{i}. {r['title']} - {r['url']}\n"
            prompt = f"{web_context}\n\n{prompt}"

    candidate_models = [
        "gemini-3.5-flash-lite",
        "gemini-flash-lite-latest",
        "gemini-flash-latest",
    ]

    api_keys_to_try = [k for k in [GEMINI_API_KEY, GEMINI_BACKUP_API_KEY] if k]

    last_error = None
    for api_key in api_keys_to_try:
        genai.configure(api_key=api_key)
        key_exhausted = False

        for model_name in candidate_models:
            if key_exhausted:
                break

            max_retries = 2
            for attempt in range(max_retries):
                try:
                    model = genai.GenerativeModel(model_name)
                    response = model.generate_content(
                        prompt,
                        generation_config=genai.types.GenerationConfig(
                            max_output_tokens=2000,
                            temperature=temp
                        )
                    )
                    try:
                        answer = response.text.strip() if response.text else None
                    except (IndexError, ValueError):
                        answer = None
                        if response.candidates:
                            for part in response.candidates[0].content.parts:
                                if hasattr(part, 'text') and part.text:
                                    answer = part.text.strip()
                                    break
                    if answer:
                        return answer
                    logger.warning(f"⚠️ {model_name} returned empty response, trying next...")
                    break
                except Exception as e:
                    error_str = str(e)
                    last_error = error_str

                    if _is_model_not_found(error_str):
                        logger.warning(f"⚠️ Model {model_name} not available, skipping")
                        break

                    if _is_rate_limit_error(error_str):
                        if attempt < max_retries - 1:
                            backoff = 2 ** (attempt + 1)
                            logger.warning(f"⚠️ {model_name} rate limited, waiting {backoff}s...")
                            _time.sleep(backoff)
                            continue
                        else:
                            logger.warning(f"⚠️ {model_name} rate limited after retries, trying next...")
                            if "quota" in error_str.lower() or "daily" in error_str.lower():
                                key_exhausted = True
                            break

                    if _is_overload_error(error_str):
                        if attempt < max_retries - 1:
                            _time.sleep(1)
                            continue
                        else:
                            logger.warning(f"⚠️ {model_name} overloaded, trying next...")
                            break

                    logger.error(f"Gemini API error ({model_name}): {error_str}")
                    return f"Error: {error_str}"

    logger.error(f"All Gemini models/keys exhausted. Last error: {last_error}")
    return ("⚠️ All available AI models have hit their request limits. "
            "Please try again in a minute or two.")


def extract_ppt_content(filepath):
    """Extract PowerPoint content using python-pptx and Gemini for analysis"""
    try:
        if not os.path.exists(filepath):
            return f"Error: File not found at {filepath}"
        
        file_size = os.path.getsize(filepath)
        if file_size > 200 * 1024 * 1024:  # 200MB limit
            return "Error: File too large (>50MB)"
        
        logger.info(f"📄 Processing PowerPoint file: {filepath} ({file_size} bytes)")
        
        try:
            prs = Presentation(filepath)
        except Exception as e:
            return f"Error: Unable to open PowerPoint file - {str(e)}"
        
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i+1} ---")
            
            text_found = False
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        text_found = True
                    
                    if hasattr(shape, "table"):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_text.extend(table_text)
                            text_found = True
                            
                except Exception as e:
                    logger.warning(f"Error processing shape on slide {i+1}: {str(e)}")
                    continue
            
            if not text_found:
                slide_text.append("[No text content found in this slide]")
                
            slides_content.append("\n".join(slide_text))
        
        if not slides_content:
            return "Error: No slides found in PowerPoint file"

        full_content = "\n\n".join(slides_content)
        logger.info(f"✓ PowerPoint content extracted ({len(full_content)} chars)")
        return full_content
        
    except Exception as e:
        logger.error(f"Error extracting PowerPoint content: {str(e)}")
        return f"Error extracting PowerPoint content: {str(e)}"

def process_document(filepath, filename):
    ext = filename.rsplit('.', 1)[1].lower()
    
    try:
        if ext in ['pptx', 'ppt']:
            content = extract_ppt_content(filepath)
            if content.startswith("Error"):
                return None, content
            
            from langchain_core.documents import Document
            
            chunks = []
            sections = content.split('--- Slide')
            for i, section in enumerate(sections):
                if section.strip():
                    if i > 0:
                        section = '--- Slide' + section
                    
                    if len(section) > 2000:
                        sub_chunks = section.split('\n\n')
                        for j, sub_chunk in enumerate(sub_chunks):
                            if sub_chunk.strip():
                                chunks.append(sub_chunk.strip())
                    else:
                        chunks.append(section.strip())
            
            pages = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    doc = Document(
                        page_content=chunk.strip(),
                        metadata={
                            "page": i + 1, 
                            "source": filename, 
                            "type": "presentation", 
                            "file_type": ext,
                            "chunk_id": i
                        }
                    )
                    pages.append(doc)
            
            logger.info(f"✓ PowerPoint processed: {len(pages)} chunks created")
            return pages, None
            
        elif ext in LOADERS:
            loader = LOADERS[ext](filepath, encoding='utf-8') if ext == 'txt' else LOADERS[ext](filepath)
            pages = loader.load()
            if not pages:
                return None, "No content extracted from document"
            
            for i, doc in enumerate(pages):
                doc.metadata = {"page": i + 1, "source": filename, "type": "document", "file_type": ext}
            return pages, None
        else:
            return None, f"Unsupported format: {ext}"
            
    except Exception as e:
        logger.error(f"Error processing {filename}: {str(e)}")
        return None, f"Error processing {filename}: {str(e)}"

def get_document_context(question, k=5):
    try:
        docs = hybrid_search(question, k=k)
        if not docs:
            return None, []
        
        context_parts, sources = [], set()
        for i, doc in enumerate(docs):
            context_parts.append(f"--- Section {i+1} ---")
            context_parts.append(doc.page_content.strip())
            if hasattr(doc, 'metadata') and 'source' in doc.metadata:
                sources.add(doc.metadata['source'])
        
        return "\n\n".join(context_parts), list(sources)
    except Exception as e:
        logger.error(f"Error getting document context: {str(e)}")
        return None, []

def extract_numeric_from_text(series):
    """Try to pull a representative number out of text values like '5-15 seconds',
    '85-90%', '$1,200', etc. Returns the average when a range is found, otherwise
    the single number. Returns None if nothing numeric can be extracted."""
    def parse_value(val):
        if pd.isna(val):
            return np.nan
        text = str(val)
        numbers = re.findall(r'\d+\.?\d*', text.replace(',', ''))
        if not numbers:
            return np.nan
        numbers = [float(n) for n in numbers]
        return sum(numbers) / len(numbers)  # average if it's a range like "5-15"

    converted = series.apply(parse_value)
    # Only treat this as usable if most values actually converted
    if converted.notna().sum() >= max(1, int(len(series) * 0.5)):
        return converted
    return None

def build_data_summary(df, filename=None, max_rows_for_full_dump=100):
    """Builds a compact, token-efficient summary of a dataframe for sending to Gemini.
    Cached per filename so repeated questions don't recompute."""
    if filename and filename in _data_summary_cache:
        return _data_summary_cache[filename]

    if len(df) <= max_rows_for_full_dump:
        result = df.to_string()
    else:
        parts = [f"Dataset has {len(df)} rows and {len(df.columns)} columns (showing summary).\n"]

        numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()

        if numeric_cols:
            parts.append("Numeric column statistics:")
            parts.append(df[numeric_cols].describe().to_string())
            parts.append("")

        for col in categorical_cols:
            unique_count = df[col].nunique()
            if unique_count <= 30:
                parts.append(f"Value counts for '{col}':")
                parts.append(df[col].value_counts().head(15).to_string())
                parts.append("")
            else:
                parts.append(f"'{col}': {unique_count} unique values (top 10):")
                parts.append(df[col].value_counts().head(10).to_string())
                parts.append("")

        parts.append(f"Sample of {min(20, len(df))} rows:")
        parts.append(df.head(20).to_string())
        result = "\n".join(parts)

    if filename:
        _data_summary_cache[filename] = result
    return result

def get_numeric_columns(df):
    return df.select_dtypes(include=[np.number]).columns.tolist()

def get_extractable_numeric_columns(df):
    """Find text columns that contain numbers we can extract (e.g. '5-15 sec', '85%')."""
    extractable = {}
    for col in df.select_dtypes(include=['object', 'category']).columns:
        converted = extract_numeric_from_text(df[col])
        if converted is not None:
            extractable[col] = converted
    return extractable

def find_mentioned_columns(question, columns):
    """Detects which dataframe columns the user actually named in their question.
    Handles both exact phrases ('pie chart of Country' -> Country) and partial
    word matches ('based on location' -> Customer Location), since users rarely
    type a multi-word column name in full. Scored by how much of the column
    name's words are present, so stronger/more-complete matches rank first."""
    q_lower = question.lower()
    q_words = set(re.findall(r'\w+', q_lower))

    scored = []
    for col in columns:
        col_lower = col.lower()
        if col_lower in q_lower:
            scored.append((col, 1.0 + len(col_lower)))  # exact phrase match, ranked highest
            continue
        col_words = set(re.findall(r'\w+', col_lower))
        # ignore tiny/common words so 'ID' or 'Date' alone don't match too eagerly
        meaningful_words = {w for w in col_words if len(w) > 2}
        if not meaningful_words:
            continue
        overlap = meaningful_words & q_words
        if overlap:
            score = len(overlap) / len(meaningful_words)
            if score >= 0.5:  # at least half the column's meaningful words were said
                scored.append((col, score))

    scored.sort(key=lambda x: x[1], reverse=True)
    return [col for col, _ in scored]

def suggest_chart_columns(df, question):
    numeric_cols = get_numeric_columns(df)
    categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
    question_lower = question.lower()

    # Honor columns the user explicitly named, instead of always auto-picking
    mentioned = find_mentioned_columns(question, df.columns.tolist())
    mentioned_categorical = [c for c in mentioned if c in categorical_cols]
    mentioned_numeric_pre = [c for c in mentioned if c in numeric_cols]

    def pick(preferred_list, fallback_list, index=0):
        if len(preferred_list) > index:
            return preferred_list[index]
        return fallback_list[index] if len(fallback_list) > index else None

    # "Number of X" / "count of X" / "how many X" means COUNT ROWS per category,
    # not sum some unrelated numeric column. This works even with zero numeric
    # columns, since it's just counting occurrences.
    count_intent_phrases = ['number of', 'count of', 'how many', 'total number',
                             'orders count', 'frequency of', 'how much orders']
    is_count_request = any(phrase in question_lower for phrase in count_intent_phrases)

    if is_count_request and categorical_cols:
        x_col = pick(mentioned_categorical, categorical_cols)
        chart_type = 'pie' if 'pie' in question_lower else 'bar'
        return chart_type, x_col, '__COUNT__', f"{chart_type.capitalize()} chart of order count by {x_col}"

    # If no true numeric columns, try extracting numbers from text columns
    extractable = {}
    if not numeric_cols:
        extractable = get_extractable_numeric_columns(df)
        if extractable:
            # Use the extracted numeric versions; remaining text columns stay categorical
            numeric_cols = list(extractable.keys())
            categorical_cols = [c for c in categorical_cols if c not in extractable]

    if not numeric_cols:
        return None, None, None, "No numeric columns found"

    mentioned_numeric = [c for c in mentioned if c in numeric_cols]

    if 'pie' in question_lower and categorical_cols:
        x_col = pick(mentioned_categorical, categorical_cols)
        y_col = pick(mentioned_numeric, numeric_cols)
        label = f"Pie chart of {y_col} by {x_col}" + (" (columns you specified)" if mentioned else "")
        return 'pie', x_col, y_col, label
    elif ('line' in question_lower or 'scatter' in question_lower) and mentioned_categorical and not mentioned_numeric:
        # User asked for line/scatter but only named a categorical column (e.g. "line
        # chart based on payment methods") — those chart types need two numeric axes,
        # which the mentioned column can't provide. A bar chart of that category
        # against a numeric measure actually answers what they're asking about,
        # instead of silently ignoring the column they named.
        x_col = mentioned_categorical[0]
        y_col = numeric_cols[0]
        return 'bar', x_col, y_col, (f"Bar chart of {y_col} by {x_col} "
                                      f"(a line/scatter chart needs two numeric measures, "
                                      f"so this shows {x_col} instead, which is what you asked about)")
    elif 'line' in question_lower and len(numeric_cols) >= 2:
        x_col = pick(mentioned_numeric, numeric_cols, 0)
        y_col = pick(mentioned_numeric, numeric_cols, 1) if len(mentioned_numeric) < 2 else mentioned_numeric[1]
        return 'line', x_col, y_col, f"Line chart of {y_col} over {x_col}"
    elif 'scatter' in question_lower and len(numeric_cols) >= 2:
        x_col = pick(mentioned_numeric, numeric_cols, 0)
        y_col = pick(mentioned_numeric, numeric_cols, 1) if len(mentioned_numeric) < 2 else mentioned_numeric[1]
        return 'scatter', x_col, y_col, f"Scatter plot of {y_col} vs {x_col}"
    elif 'histogram' in question_lower or 'distribution' in question_lower:
        y_col = pick(mentioned_numeric, numeric_cols)
        return 'histogram', None, y_col, f"Histogram showing distribution of {y_col}"
    elif categorical_cols and numeric_cols:
        x_col = pick(mentioned_categorical, categorical_cols)
        y_col = pick(mentioned_numeric, numeric_cols)
        return 'bar', x_col, y_col, f"Bar chart of {y_col} by {x_col}" + (" (columns you specified)" if mentioned else "")
    elif len(numeric_cols) >= 2:
        x_col = pick(mentioned_numeric, numeric_cols, 0)
        y_col = pick(mentioned_numeric, numeric_cols, 1) if len(mentioned_numeric) < 2 else mentioned_numeric[1]
        return 'scatter', x_col, y_col, f"Scatter plot of {y_col} vs {x_col}"
    else:
        return 'histogram', None, numeric_cols[0], "Distribution histogram"

def generate_smart_chart(df, filename, question):
    """Generate chart with improved error handling and better visualization"""
    chart_type, x_col, y_col, reasoning = suggest_chart_columns(df, question)
    if not chart_type:
        return None, reasoning

    df = df.copy()

    # '__COUNT__' is a sentinel meaning "count rows per category" (e.g. "number of
    # orders per location") rather than summing an actual numeric column. Seeding
    # a column of 1s lets the existing sum-based grouping logic below produce
    # counts naturally, with no special-casing needed in the plotting code.
    if y_col == '__COUNT__':
        df['__COUNT__'] = 1

    # If y_col (or x_col for histograms) isn't a real numeric dtype, it must be one of our
    # extracted text-based numeric columns — replace it with the parsed numeric series.
    for col in [x_col, y_col]:
        if col and col in df.columns and df[col].dtype == 'object':
            converted = extract_numeric_from_text(df[col])
            if converted is not None:
                df[col] = converted

    try:
        chart_id = uuid.uuid4().hex
        
        # Set up the plot with better styling
        plt.figure(figsize=(12, 8))
        plt.style.use('dark_background')
        
        # Color palette
        colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
        
        if chart_type == 'bar' and x_col and y_col:
            if df[x_col].dtype == 'object':
                # Group and aggregate data
                grouped = df.groupby(x_col)[y_col].sum().sort_values(ascending=False).head(10)
                bars = plt.bar(range(len(grouped)), grouped.values, color=colors[0], alpha=0.8)
                plt.xticks(range(len(grouped)), grouped.index, rotation=45, ha='right')
                
                # Add value labels on bars
                for bar in bars:
                    height = bar.get_height()
                    plt.text(bar.get_x() + bar.get_width()/2., height,
                            f'{height:.1f}', ha='center', va='bottom')
            else:
                plt.bar(df[x_col], df[y_col], color=colors[0], alpha=0.8)
            
            plt.xlabel(x_col, fontsize=12)
            plt.ylabel(y_col, fontsize=12)
            plt.title(f'{y_col} by {x_col}', fontsize=14, fontweight='bold')
            
        elif chart_type == 'scatter' and x_col and y_col:
            plt.scatter(df[x_col], df[y_col], alpha=0.7, color=colors[0], s=60)
            plt.xlabel(x_col, fontsize=12)
            plt.ylabel(y_col, fontsize=12)
            plt.title(f'{y_col} vs {x_col}', fontsize=14, fontweight='bold')
            
            # Add trend line
            z = np.polyfit(df[x_col].dropna(), df[y_col].dropna(), 1)
            p = np.poly1d(z)
            plt.plot(df[x_col], p(df[x_col]), "--", color=colors[1], alpha=0.8)
            
        elif chart_type == 'histogram' and y_col:
            plt.hist(df[y_col].dropna(), bins=20, color=colors[0], alpha=0.7, edgecolor='white')
            plt.xlabel(y_col, fontsize=12)
            plt.ylabel('Frequency', fontsize=12)
            plt.title(f'Distribution of {y_col}', fontsize=14, fontweight='bold')
            
        elif chart_type == 'pie' and x_col and y_col:
            if df[x_col].dtype == 'object':
                pie_data = df.groupby(x_col)[y_col].sum().sort_values(ascending=False).head(8)  # Actual top 8 categories
                plt.pie(pie_data.values, labels=pie_data.index, autopct='%1.1f%%', 
                       colors=colors, startangle=90)
                plt.title(f'{y_col} Distribution by {x_col}', fontsize=14, fontweight='bold')
                
        elif chart_type == 'line' and x_col and y_col:
            plt.plot(df[x_col], df[y_col], marker='o', linewidth=2, markersize=6, color=colors[0])
            plt.xlabel(x_col, fontsize=12)
            plt.ylabel(y_col, fontsize=12)
            plt.title(f'{y_col} over {x_col}', fontsize=14, fontweight='bold')
            plt.grid(True, alpha=0.3)
        
        # Improve layout and save
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.15)
        
        # Ensure the plots directory exists
        os.makedirs(PLOTS_FOLDER, exist_ok=True)
        
        chart_path = os.path.join(PLOTS_FOLDER, f"{chart_id}.png")
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='#1a1f2e', 
                   edgecolor='none', transparent=False)
        plt.close()
        
        # Verify file was created
        if os.path.exists(chart_path):
            logger.info(f"✓ Chart saved successfully: {chart_path}")
            # Return the full URL for the frontend
            chart_url = f"/static/plots/{chart_id}.png"
            return chart_url, f"Generated {chart_type} chart: {reasoning}"
        else:
            logger.error("❌ Chart file was not created")
            return None, "Failed to save chart file"
        
    except Exception as e:
        plt.close()  # Ensure plot is closed on error
        logger.error(f"Error generating chart: {str(e)}")
        return None, f"Error generating chart: {str(e)}"

FRONTEND_DIR = os.path.dirname(os.path.abspath(__file__))

# ─── AUTH ROUTES ─────────────────────────────────────────────────────────────

@app.route('/auth/signup', methods=['POST'])
def auth_signup():
    if not pg_engine:
        return jsonify({'error': 'Database not configured'}), 503
    data = request.get_json() or {}
    fullname = data.get('fullname', '').strip()
    username = data.get('username', '').strip().lower()
    email = data.get('email', '').strip().lower()
    password = data.get('password', '')

    if not email or not password:
        return jsonify({'error': 'Email and password are required'}), 400
    if len(password) < 8:
        return jsonify({'error': 'Password must be at least 8 characters'}), 400
    if not re.match(r'^[^\s@]+@[^\s@]+\.[^\s@]+$', email):
        return jsonify({'error': 'Invalid email address'}), 400

    try:
        with pg_engine.connect() as conn:
            existing = conn.execute(text("SELECT id FROM users WHERE email = :email"), {'email': email}).fetchone()
            if existing:
                return jsonify({'error': 'Email already registered'}), 409
            if username:
                existing = conn.execute(text("SELECT id FROM users WHERE username = :username"), {'username': username}).fetchone()
                if existing:
                    return jsonify({'error': 'Username already taken'}), 409

            password_hash = generate_password_hash(password)
            row = conn.execute(text("""
                INSERT INTO users (fullname, username, email, password_hash, provider)
                VALUES (:fullname, :username, :email, :password_hash, 'email')
                RETURNING id, email, fullname, username
            """), {'fullname': fullname, 'username': username, 'email': email, 'password_hash': password_hash}).fetchone()
            conn.commit()
        logger.info(f"✓ New user registered: {email}")
        return jsonify({'success': True, 'message': 'Account created successfully'}), 201
    except Exception as e:
        logger.error(f"Signup error: {str(e)}")
        return jsonify({'error': 'Registration failed. Please try again.'}), 500


@app.route('/auth/login', methods=['POST'])
def auth_login():
    if not pg_engine:
        return jsonify({'error': 'Database not configured'}), 503
    data = request.get_json() or {}
    email = data.get('email', '').strip().lower()
    password = data.get('password', '')

    if not email or not password:
        return jsonify({'error': 'Email and password are required'}), 400

    try:
        with pg_engine.connect() as conn:
            user = conn.execute(text(
                "SELECT id, email, fullname, username, password_hash FROM users WHERE email = :email"
            ), {'email': email}).fetchone()

        if not user:
            return jsonify({'error': 'Invalid email or password'}), 401
        if not user[4] or not check_password_hash(user[4], password):
            return jsonify({'error': 'Invalid email or password'}), 401

        token = pyjwt.encode({
            'sub': str(user[0]),
            'id': str(user[0]),
            'email': user[1],
            'name': user[2] or '',
            'username': user[3] or '',
            'exp': datetime.now(timezone.utc) + timedelta(days=7)
        }, JWT_SECRET, algorithm='HS256')

        initials = ''
        if user[2]:
            initials = ''.join(w[0] for w in user[2].split() if w)[:2].upper()
        else:
            initials = email[:2].upper()

        logger.info(f"✓ User logged in: {email}")
        return jsonify({
            'success': True,
            'token': token,
            'user': {
                'id': str(user[0]),
                'email': user[1],
                'name': user[2] or '',
                'username': user[3] or '',
                'initials': initials
            }
        })
    except Exception as e:
        logger.error(f"Login error: {str(e)}")
        return jsonify({'error': 'Login failed. Please try again.'}), 500


@app.route('/auth/verify', methods=['GET'])
def auth_verify():
    auth_header = request.headers.get('Authorization', '')
    token = auth_header[7:] if auth_header.startswith('Bearer ') else None
    if not token:
        return jsonify({'valid': False, 'error': 'No token'}), 401
    try:
        payload = pyjwt.decode(token, JWT_SECRET, algorithms=['HS256'])
        return jsonify({
            'valid': True,
            'user': {
                'id': payload.get('id'),
                'name': payload.get('name'),
                'email': payload.get('email'),
                'username': payload.get('username')
            }
        })
    except pyjwt.ExpiredSignatureError:
        return jsonify({'valid': False, 'error': 'Token expired'}), 401
    except Exception:
        return jsonify({'valid': False, 'error': 'Invalid token'}), 401


@app.route('/auth/logout', methods=['POST'])
def auth_logout():
    return jsonify({'success': True})


@app.route('/auth/google')
def auth_google():
    if not GOOGLE_CLIENT_ID or not GOOGLE_CLIENT_SECRET:
        return jsonify({'error': 'Google OAuth not configured. Set GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET in .env'}), 503
    from_page = request.args.get('from', 'login')
    state = from_page
    params = {
        'client_id': GOOGLE_CLIENT_ID,
        'redirect_uri': GOOGLE_REDIRECT_URI,
        'response_type': 'code',
        'scope': 'openid email profile',
        'access_type': 'offline',
        'prompt': 'select_account',
        'state': state,
    }
    url = 'https://accounts.google.com/o/oauth2/v2/auth?' + '&'.join(f'{k}={requests.utils.quote(str(v))}' for k, v in params.items())
    return redirect(url)


@app.route('/auth/google/callback')
def auth_google_callback():
    code = request.args.get('code')
    state = request.args.get('state', 'login')
    if not code:
        return redirect(f'/login.html?error=auth_failed')
    try:
        token_resp = requests.post('https://oauth2.googleapis.com/token', data={
            'code': code,
            'client_id': GOOGLE_CLIENT_ID,
            'client_secret': GOOGLE_CLIENT_SECRET,
            'redirect_uri': GOOGLE_REDIRECT_URI,
            'grant_type': 'authorization_code',
        }, timeout=10)
        tokens = token_resp.json()
        access_token = tokens.get('access_token')
        if not access_token:
            return redirect(f'/login.html?error=auth_failed')

        userinfo_resp = requests.get('https://www.googleapis.com/oauth2/v2/userinfo',
                                     headers={'Authorization': f'Bearer {access_token}'}, timeout=10)
        userinfo = userinfo_resp.json()
        g_email = userinfo.get('email', '')
        g_name = userinfo.get('name', '')
        g_picture = userinfo.get('picture', '')
        g_id = userinfo.get('id', '')

        with pg_engine.connect() as conn:
            existing = conn.execute(text("SELECT id, fullname FROM users WHERE email = :email"), {'email': g_email}).fetchone()
            if existing:
                user_id = str(existing[0])
                user_name = existing[1] or g_name
            else:
                row = conn.execute(text("""
                    INSERT INTO users (fullname, email, picture, provider)
                    VALUES (:name, :email, :picture, 'google')
                    RETURNING id
                """), {'name': g_name, 'email': g_email, 'picture': g_picture}).fetchone()
                conn.commit()
                user_id = str(row[0])
                user_name = g_name

        initials = ''.join(w[0] for w in g_name.split() if w)[:2].upper() if g_name else g_email[:2].upper()
        token = pyjwt.encode({
            'sub': user_id, 'id': user_id,
            'email': g_email, 'name': user_name,
            'picture': g_picture, 'initials': initials,
            'exp': datetime.now(timezone.utc) + timedelta(days=7)
        }, JWT_SECRET, algorithm='HS256')

        return_page = 'signup' if state == 'signup' else 'login'
        return redirect(f'/{return_page}.html?token={token}')
    except Exception as e:
        logger.error(f"Google OAuth error: {str(e)}")
        return redirect(f'/login.html?error=auth_failed')


# ─── PAGE ROUTES ─────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return send_from_directory(FRONTEND_DIR, 'login.html')

@app.route('/login.html')
def serve_login():
    return send_from_directory(FRONTEND_DIR, 'login.html')

@app.route('/signup.html')
def serve_signup():
    return send_from_directory(FRONTEND_DIR, 'signup.html')

@app.route('/docurag.html')
def serve_docurag():
    return send_from_directory(FRONTEND_DIR, 'docurag.html')

@app.route('/health')
def health_check():
    logger.info("Health check received")
    return jsonify({
        'status': 'healthy',
        'gemini_api': bool(GEMINI_API_KEY),
        #'has_documents': has_documents(),
        'has_data': has_data(),
        #'data_files': list(excel_data_store.keys()),
        #'upload_folder': UPLOAD_FOLDER,
        #'plots_folder': PLOTS_FOLDER,
        #'vector_db_exists': os.path.exists(VECTOR_DB_DIR),
        'postgres_connected': using_postgres()
    })

# ---------------------------------------------------------------------------
# Chat history endpoints (New Chat + History of chats)
# Only active when DATABASE_URL is set. If Postgres isn't configured yet,
# these return a clear error rather than crashing, so the rest of the app
# keeps working while you finish setting up Supabase.
# ---------------------------------------------------------------------------

def require_postgres():
    if not using_postgres():
        return jsonify({'error': 'Chat history requires Postgres. Set DATABASE_URL and restart the server.'}), 503
    return None

@app.route('/api/chats', methods=['GET'])
def list_chats():
    """Returns all chat sessions, most recent first — for the sidebar history list."""
    err = require_postgres()
    if err:
        return err
    try:
        with pg_engine.connect() as conn:
            rows = conn.execute(text(
                "SELECT id, title, created_at FROM chat_sessions ORDER BY created_at DESC"
            )).fetchall()
        return jsonify({
            'chats': [{'id': str(r[0]), 'title': r[1], 'created_at': r[2].isoformat()} for r in rows]
        })
    except Exception as e:
        logger.error(f"Error listing chats: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/chats', methods=['POST'])
def create_chat():
    """Creates a new chat session — for the 'New Chat' button."""
    err = require_postgres()
    if err:
        return err
    try:
        data = request.get_json() or {}
        title = data.get('title', 'New Chat')
        with pg_engine.connect() as conn:
            row = conn.execute(text(
                "INSERT INTO chat_sessions (title) VALUES (:title) RETURNING id, title, created_at"
            ), {'title': title}).fetchone()
            conn.commit()
        return jsonify({'id': str(row[0]), 'title': row[1], 'created_at': row[2].isoformat()})
    except Exception as e:
        logger.error(f"Error creating chat: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/chats/<chat_id>', methods=['DELETE'])
def delete_chat(chat_id):
    err = require_postgres()
    if err:
        return err
    try:
        with pg_engine.connect() as conn:
            conn.execute(text("DELETE FROM chat_sessions WHERE id = :id"), {'id': chat_id})
            conn.commit()
        return jsonify({'message': 'Chat deleted'})
    except Exception as e:
        logger.error(f"Error deleting chat: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/chats/<chat_id>/messages', methods=['GET'])
def get_chat_messages(chat_id):
    """Returns full message history for a chat session — for loading history."""
    err = require_postgres()
    if err:
        return err
    try:
        with pg_engine.connect() as conn:
            rows = conn.execute(text(
                "SELECT role, content, mode, sources, chart_url, created_at "
                "FROM chat_messages WHERE session_id = :sid ORDER BY created_at ASC"
            ), {'sid': chat_id}).fetchall()
        return jsonify({
            'messages': [{
                'role': r[0], 'content': r[1], 'mode': r[2],
                'sources': json.loads(r[3]) if r[3] else [],
                'chart_url': r[4], 'created_at': r[5].isoformat()
            } for r in rows]
        })
    except Exception as e:
        logger.error(f"Error fetching messages: {str(e)}")
        return jsonify({'error': str(e)}), 500

def get_conversation_history(session_id, limit=10):
    """Fetches recent conversation history for context continuity. Returns formatted string."""
    if not using_postgres() or not session_id:
        return ""
    try:
        with pg_engine.connect() as conn:
            rows = conn.execute(text(
                "SELECT role, content FROM chat_messages "
                "WHERE session_id = :sid ORDER BY created_at DESC LIMIT :lim"
            ), {'sid': session_id, 'lim': limit}).fetchall()
        if not rows:
            return ""
        rows = list(reversed(rows))
        history_lines = []
        for role, content in rows:
            prefix = "User" if role == "user" else "Assistant"
            truncated = content[:500] if len(content) > 500 else content
            history_lines.append(f"{prefix}: {truncated}")
        return "\n".join(history_lines)
    except Exception as e:
        logger.error(f"Error fetching conversation history: {str(e)}")
        return ""


def save_message(session_id, role, content, mode=None, sources=None, chart_url=None):
    """Persists a single chat message. Silently no-ops if Postgres isn't configured."""
    if not using_postgres() or not session_id:
        return
    try:
        with pg_engine.connect() as conn:
            conn.execute(text(
                "INSERT INTO chat_messages (session_id, role, content, mode, sources, chart_url) "
                "VALUES (:sid, :role, :content, :mode, :sources, :chart_url)"
            ), {
                'sid': session_id, 'role': role, 'content': content, 'mode': mode,
                'sources': json.dumps(sources) if sources else None, 'chart_url': chart_url
            })
            conn.commit()
    except Exception as e:
        logger.error(f"Error saving message: {str(e)}")

@app.route('/upload', methods=['POST'])
def upload():
    global last_uploaded_type, _documents_indexed
    logger.info("Upload request received")

    if 'file' not in request.files:
        return jsonify({'error': 'No file part in request'}), 400

    file = request.files['file']
    if not file or file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    logger.info(f"Filename received: {file.filename}")

    if not allowed_file(file.filename):
        return jsonify({'error': 'Unsupported file type'}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)

    try:
        file.save(filepath)
        logger.info(f"File saved: {filepath}")
    except Exception as e:
        logger.error(f"Failed to save file: {str(e)}")
        return jsonify({'error': f'Failed to save file: {str(e)}'}), 500

    ext = filename.rsplit('.', 1)[1].lower()

    # Handle data files
    if ext in ['xls', 'xlsx', 'csv']:
        try:
            df = pd.read_csv(filepath) if ext == 'csv' else pd.read_excel(filepath)
            excel_data_store[filename] = df
            _data_summary_cache.pop(filename, None)
            last_uploaded_type = 'data'
            logger.info(f"📊 Data file processed: {filename} - {len(df)} rows, {len(df.columns)} columns")
            logger.info(f"📊 Columns: {df.columns.tolist()}")
            logger.info(f"📊 Data types: {df.dtypes.to_dict()}")
            
            return jsonify({
                'success': True,
                'message': f'File uploaded successfully',
                'filename': filename,
                'columns': df.columns.tolist(),
                'shape': df.shape,
                'data_types': df.dtypes.astype(str).to_dict(),
                'sample_data': df.head(3).to_dict()
            })
        except Exception as e:
            logger.error(f"Failed to process data file: {str(e)}")
            return jsonify({'error': f'Failed to process data file: {str(e)}'}), 500
    
    # Handle document files
    else:
        logger.info("Document processing started")

        pages, error = process_document(filepath, filename)
        if error:
            logger.error(f"Document processing error: {error}")
            return jsonify({'error': error}), 500

        try:
            total_chars = sum(len(p.page_content) for p in pages)
            chunk_size = 1500 if total_chars > 50000 else 1000
            chunk_overlap = 150 if chunk_size == 1500 else 100

            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
            )
            texts = text_splitter.split_documents(pages)

            BATCH_SIZE = 200
            for i in range(0, len(texts), BATCH_SIZE):
                batch = texts[i:i + BATCH_SIZE]
                index_documents(batch)

            last_uploaded_type = 'document'
            _documents_indexed = True

            logger.info(f"Document processing completed: {len(pages)} pages into {len(texts)} chunks")
            return jsonify({
                'success': True,
                'message': f'File uploaded successfully',
                'filename': filename,
                'pages': len(pages),
                'chunks': len(texts),
                'file_type': ext
            })
        except Exception as e:
            logger.error(f"Failed to index document: {str(e)}")
            return jsonify({'error': f'Failed to index document: {str(e)}'}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.get_json()
    question = data.get('question', '')
    rag_mode = data.get('rag_mode', True)  # New parameter to control RAG vs Web mode
    session_id = data.get('session_id')  # optional — enables chat history persistence
    
    if not question:
        return jsonify({'error': 'No question provided'}), 400
    
    logger.info(f"❓ Question received: {question}")
    logger.info(f"🔄 Mode: {'RAG' if rag_mode else 'Web Search'}")
    
    # Fetch conversation history BEFORE saving current message (so it contains only prior Q&A)
    conversation_history = get_conversation_history(session_id, limit=10)

    # Save the user's question (if a chat session is active)
    save_message(session_id, 'user', question)

    def respond(response_data, status=200):
        """Wraps a jsonify response and persists the assistant's answer to history."""
        save_message(
            session_id, 'assistant',
            response_data.get('answer', ''),
            mode=response_data.get('mode'),
            sources=response_data.get('sources'),
            chart_url=response_data.get('chart_url')
        )
        return jsonify(response_data), status
    
    is_graph_req = is_graph_request(question)
    is_doc_question = is_document_specific_question(question)
    
    try:
        # RAG Mode - Use uploaded files and documents
        if rag_mode:
            logger.info("📚 RAG Mode: Processing with uploaded files...")
            
            # Handle graph requests with priority
            if is_graph_req and has_data():
                logger.info("📊 Processing graph request...")
                filename, df = list(excel_data_store.items())[-1]
                logger.info(f"📊 Using dataset: {filename} with shape {df.shape}")
                
                chart_url, chart_info = generate_smart_chart(df, filename, question)
                
                if chart_url:
                    logger.info(f"✅ Chart generated successfully: {chart_url}")
                    
                    history_block = f"\nCONVERSATION HISTORY:\n{conversation_history}\n" if conversation_history else ""
                    analysis_prompt = f"""I've created a chart based on: "{question}"
Dataset: {filename} ({df.shape[0]} rows × {df.shape[1]} columns)
Chart: {chart_info}
Columns: {', '.join(df.columns.tolist())}
Sample data:
{df.head(3).to_string()}
{history_block}
Provide insightful analysis of this visualization and explain what the chart shows."""
                    
                    return respond({
                        'answer': ask_gemini(analysis_prompt, temp=0.5),
                        'chart_url': chart_url,
                        'chart_info': chart_info,
                        'mode': 'rag_graph_generated'
                    })
                else:
                    logger.error(f"❌ Chart generation failed: {chart_info}")
                    return respond({
                        'answer': f"I couldn't generate the chart. {chart_info}. Please check your data format and try again.",
                        'mode': 'rag_graph_failed',
                        'error': chart_info
                    })
            
            # Handle document questions
            elif last_uploaded_type == 'data' and has_data():
                filename, df = list(excel_data_store.items())[-1]
                logger.info(f"📊 Answering data question using dataset: {filename}")

                history_block = f"\nCONVERSATION HISTORY:\n{conversation_history}\n" if conversation_history else ""
                data_prompt = f"""Answer the following question using this dataset:

Dataset: {filename} ({df.shape[0]} rows × {df.shape[1]} columns)
Columns: {', '.join(df.columns.tolist())}

Data:
{build_data_summary(df, filename=filename)}
{history_block}
QUESTION: {question}

Use the conversation history above (if any) to understand follow-up questions in context. Provide a clear, direct answer based on the data above. If the data was summarized due to size, mention that your answer is based on the full dataset's statistics/sample, not every individual row. If the answer cannot be found in the data, say so clearly instead of guessing."""

                answer = ask_gemini(data_prompt, temp=0.3)
                return respond({
                    'answer': answer,
                    'mode': 'rag_data_qa',
                    'filename': filename
                })

            elif has_documents():
                context, sources = get_document_context(question, k=5)

                if context:
                    history_block = f"\nCONVERSATION HISTORY:\n{conversation_history}\n" if conversation_history else ""
                    rag_prompt = f"""Based on the following document content, answer the question comprehensively:

DOCUMENT CONTENT:
{context}
{history_block}
QUESTION: {question}

Use the conversation history above (if any) to understand follow-up questions in context. For example, if the user says "explain the second point" or "tell me more about that", refer to your previous answers and the document to give a coherent response. Provide a detailed, accurate response using ONLY the document information. If the answer cannot be found in the document, clearly state that instead of making up information."""

                    answer = ask_gemini(rag_prompt, temp=0.3)
                    return respond({
                        'answer': answer,
                        'sources': sources,
                        'mode': 'rag_document',
                        'context_found': True
                    })
                else:
                    return respond({
                        'answer': "No relevant information found in uploaded documents for your question.",
                        'mode': 'rag_no_context',
                        'context_found': False
                    })
            
            # Handle general questions about uploaded data (fallback if type wasn't tracked)
            elif has_data():
                filename, df = list(excel_data_store.items())[-1]
                logger.info(f"📊 Answering data question using dataset: {filename}")

                history_block = f"\nCONVERSATION HISTORY:\n{conversation_history}\n" if conversation_history else ""
                data_prompt = f"""Answer the following question using this dataset:

Dataset: {filename} ({df.shape[0]} rows × {df.shape[1]} columns)
Columns: {', '.join(df.columns.tolist())}

Data:
{build_data_summary(df, filename=filename)}
{history_block}
QUESTION: {question}

Use the conversation history above (if any) to understand follow-up questions in context. Provide a clear, direct answer based on the data above. If the data was summarized due to size, mention that your answer is based on the full dataset's statistics/sample, not every individual row. If the answer cannot be found in the data, say so clearly instead of guessing."""

                answer = ask_gemini(data_prompt, temp=0.3)
                return respond({
                    'answer': answer,
                    'mode': 'rag_data_qa',
                    'filename': filename
                })

            # No uploaded files in RAG mode
            else:
                return respond({
                    'answer': "RAG mode is ON but no files are uploaded. Please upload documents or data files to use RAG functionality, or switch to Web Search mode for general questions.",
                    'mode': 'rag_no_files'
                })
        
        # Web Search Mode - Use Gemini with web context
        else:
            logger.info("🌐 Web Search Mode: Processing with web search...")
            
            # Use Gemini with web search context for current information
            web_enhanced_prompt = f"""Please provide a comprehensive, up-to-date answer to this question: {question}

Use your knowledge and reasoning to provide the most accurate and current information available."""
            
            answer = ask_gemini(web_enhanced_prompt, temp=0.7, include_web_context=True, web_query=question)
            
            return respond({
                'answer': answer,
                'mode': 'web_search',
                'web_enhanced': True
            })
        
    except Exception as e:
        logger.error(f"Error processing question: {str(e)}")
        return respond({
            'answer': f"Error processing question: {str(e)}",
            'mode': 'error'
        })

@app.route('/static/plots/<filename>')
def serve_plot(filename):
    """Serve generated plot images"""
    try:
        return send_from_directory(PLOTS_FOLDER, filename)
    except Exception as e:
        logger.error(f"Error serving plot {filename}: {str(e)}")
        return jsonify({'error': 'Plot not found'}), 404

@app.route('/clear', methods=['POST'])
def clear_data():
    """Clear all uploaded data and documents"""
    try:
        # Clear vector database
        clear_vector_database()

        # Clear data store
        excel_data_store.clear()
        _data_summary_cache.clear()
        
        # Clear upload folder
        if os.path.exists(UPLOAD_FOLDER):
            for filename in os.listdir(UPLOAD_FOLDER):
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                try:
                    if os.path.isfile(file_path):
                        os.unlink(file_path)
                except Exception as e:
                    logger.error(f"Error deleting {file_path}: {str(e)}")
        
        # Clear plots folder
        if os.path.exists(PLOTS_FOLDER):
            for filename in os.listdir(PLOTS_FOLDER):
                file_path = os.path.join(PLOTS_FOLDER, filename)
                try:
                    if os.path.isfile(file_path):
                        os.unlink(file_path)
                except Exception as e:
                    logger.error(f"Error deleting {file_path}: {str(e)}")
        
        logger.info("🧹 All data cleared successfully")
        return jsonify({'message': 'All data cleared successfully'})
        
    except Exception as e:
        logger.error(f"Error clearing data: {str(e)}")
        return jsonify({'error': f'Error clearing data: {str(e)}'}), 500

@app.route('/toggle-mode', methods=['POST'])
def toggle_mode():
    """Endpoint to get current mode status"""
    data = request.get_json()
    rag_mode = data.get('rag_mode', True)
    
    # Return current status
    return jsonify({
        'rag_mode': rag_mode,
        'has_documents': has_documents(),
        'has_data': has_data(),
        'message': f"Mode set to {'RAG (Files)' if rag_mode else 'Web Search'}"
    })

@app.route('/get-data-info', methods=['GET'])
def get_data_info():
    """Get information about uploaded data files"""
    if not excel_data_store:
        return jsonify({
            'has_data': False,
            'message': 'No data files uploaded'
        })
    
    data_info = {}
    for filename, df in excel_data_store.items():
        data_info[filename] = {
            'shape': df.shape,
            'columns': df.columns.tolist(),
            'numeric_columns': get_numeric_columns(df),
            'data_types': df.dtypes.astype(str).to_dict(),
            'sample': df.head(2).to_dict()
        }
    
    return jsonify({
        'has_data': True,
        'files': data_info,
        'total_files': len(excel_data_store)
    })

@app.route('/get-document-info', methods=['GET'])
def get_document_info():
    """Get information about uploaded documents"""
    try:
        if not has_documents():
            return jsonify({
                'has_documents': False,
                'message': 'No documents uploaded'
            })
        
        db = get_vector_store()
        
        # Get a sample to understand document structure
        sample_docs = db.similarity_search("sample", k=3)
        
        doc_info = {
            'has_documents': True,
            'total_chunks': len(sample_docs) if sample_docs else 0,
            'sample_content': []
        }
        
        for i, doc in enumerate(sample_docs[:2]):
            doc_info['sample_content'].append({
                'chunk_id': i,
                'content_preview': doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content,
                'metadata': doc.metadata
            })
        
        return jsonify(doc_info)
        
    except Exception as e:
        logger.error(f"Error getting document info: {str(e)}")
        return jsonify({
            'has_documents': False,
            'error': str(e)
        })

@app.route('/analyze-data', methods=['POST'])
def analyze_data():
    """Provide statistical analysis of uploaded data"""
    if not excel_data_store:
        return jsonify({'error': 'No data files uploaded'}), 400
    
    data = request.get_json()
    filename = data.get('filename')
    
    try:
        if filename and filename in excel_data_store:
            df = excel_data_store[filename]
        else:
            # Use first available dataset
            filename, df = list(excel_data_store.items())[-1]
        
        analysis = {
            'filename': filename,
            'shape': df.shape,
            'columns': df.columns.tolist(),
            'numeric_columns': get_numeric_columns(df),
            'missing_values': df.isnull().sum().to_dict(),
            'data_types': df.dtypes.astype(str).to_dict()
        }

        # Basic statistics for numeric columns
        numeric_cols = get_numeric_columns(df)
        if numeric_cols:
            analysis['statistics'] = df[numeric_cols].describe().to_dict()

        # Comprehensive per-column analysis
        column_analysis = {}
        for col in df.columns:
            col_data = df[col]
            total = len(col_data)
            missing = int(col_data.isnull().sum())
            non_null = total - missing
            unique = int(col_data.nunique())
            dtype_str = str(col_data.dtype)

            col_info = {
                'name': col,
                'dtype': dtype_str,
                'total': total,
                'non_null': non_null,
                'missing': missing,
                'missing_pct': round(missing / total * 100, 1) if total > 0 else 0,
                'unique': unique,
                'cardinality': round(unique / non_null * 100, 1) if non_null > 0 else 0,
            }

            if col_data.dtype in ['int64', 'float64', 'int32', 'float32']:
                col_info['col_type'] = 'numeric'
                desc = col_data.describe()
                col_info['min'] = None if pd.isna(desc.get('min')) else float(desc['min'])
                col_info['max'] = None if pd.isna(desc.get('max')) else float(desc['max'])
                col_info['mean'] = None if pd.isna(desc.get('mean')) else float(desc['mean'])
                col_info['median'] = None if pd.isna(col_data.median()) else float(col_data.median())
                col_info['std'] = None if pd.isna(desc.get('std')) else float(desc['std'])
                q1 = float(desc['25%']) if not pd.isna(desc.get('25%')) else None
                q3 = float(desc['75%']) if not pd.isna(desc.get('75%')) else None
                if q1 is not None and q3 is not None:
                    iqr = q3 - q1
                    outliers = int(((col_data < (q1 - 1.5 * iqr)) | (col_data > (q3 + 1.5 * iqr))).sum())
                    col_info['outliers'] = outliers
                    col_info['q1'] = q1
                    col_info['q3'] = q3
            elif col_data.dtype == 'bool':
                col_info['col_type'] = 'boolean'
                vc = col_data.value_counts().head(2)
                col_info['top_values'] = {str(k): int(v) for k, v in vc.items()}
            elif 'datetime' in dtype_str:
                col_info['col_type'] = 'datetime'
                col_info['min'] = str(col_data.min()) if not pd.isna(col_data.min()) else None
                col_info['max'] = str(col_data.max()) if not pd.isna(col_data.max()) else None
            else:
                if unique <= 20 and non_null > 0:
                    col_info['col_type'] = 'categorical'
                else:
                    col_info['col_type'] = 'text'
                vc = col_data.value_counts().head(5)
                col_info['top_values'] = {str(k): int(v) for k, v in vc.items()}
                if col_info['col_type'] == 'text' and non_null > 0:
                    lengths = col_data.dropna().astype(str).str.len()
                    col_info['avg_length'] = round(float(lengths.mean()), 1)

            # Duplicate info
            dup_count = int(col_data.duplicated().sum())
            col_info['duplicates'] = dup_count
            col_info['dup_pct'] = round(dup_count / total * 100, 1) if total > 0 else 0

            column_analysis[col] = col_info

        analysis['column_analysis'] = column_analysis

        # Generate analysis with Gemini
        analysis_prompt = f"""Analyze this dataset:

Dataset: {filename}
Shape: {df.shape[0]} rows × {df.shape[1]} columns
Columns: {', '.join(df.columns.tolist())}
Numeric columns: {numeric_cols}
Missing values: {dict(df.isnull().sum())}

Sample data:
{df.head(5).to_string()}

Provide insights about:
1. Data quality and completeness
2. Key patterns or trends visible
3. Potential analysis opportunities
4. Recommended visualizations
5. Data preparation suggestions"""

        ai_analysis = ask_gemini(analysis_prompt, temp=0.5)
        analysis['ai_insights'] = ai_analysis
        
        return jsonify(analysis)
        
    except Exception as e:
        logger.error(f"Error analyzing data: {str(e)}")
        return jsonify({'error': f'Error analyzing data: {str(e)}'}), 500

@app.route('/suggest-questions', methods=['POST'])
def suggest_questions():
    """Suggest relevant questions based on uploaded content"""
    suggestions = []
    
    try:
        # Suggestions for data files
        if excel_data_store:
            for filename, df in excel_data_store.items():
                numeric_cols = get_numeric_columns(df)
                categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
                
                if numeric_cols and categorical_cols:
                    suggestions.extend([
                        f"Create a bar chart of {numeric_cols[0]} by {categorical_cols[0]}",
                        f"Show me the distribution of {numeric_cols[0]}",
                        f"Create a scatter plot of {numeric_cols[0]} vs {numeric_cols[1] if len(numeric_cols) > 1 else numeric_cols[0]}",
                        f"Analyze the relationship between {categorical_cols[0]} and {numeric_cols[0]}",
                        f"What are the key insights from the {filename} data?"
                    ])
                elif numeric_cols:
                    suggestions.extend([
                        f"Show me the distribution of {numeric_cols[0]}",
                        f"Create a histogram of {numeric_cols[0]}",
                        f"What are the statistics for {numeric_cols[0]}?"
                    ])
        
        # Suggestions for documents
        if has_documents():
            suggestions.extend([
                "Summarize the main points from the document",
                "What are the key findings mentioned?",
                "Explain the methodology described in the document",
                "What recommendations are provided?",
                "Analyze the conclusions presented"
            ])
        
        # General suggestions if no files
        if not excel_data_store and not has_documents():
            suggestions.extend([
                "What is machine learning?",
                "Explain artificial intelligence concepts",
                "How does data analysis work?",
                "What are the latest trends in technology?",
                "Tell me about data visualization best practices"
            ])
        
        return jsonify({
            'suggestions': suggestions[:8],  # Limit to 8 suggestions
            'has_data': len(excel_data_store) > 0,
            'has_documents': has_documents()
        })
        
    except Exception as e:
        logger.error(f"Error generating suggestions: {str(e)}")
        return jsonify({
            'suggestions': ["Ask me anything!"],
            'error': str(e)
        })

# Error handlers
@app.errorhandler(404)
def not_found_error(error):
    return jsonify({'error': 'Endpoint not found'}), 404

@app.errorhandler(500)
def internal_error(error):
    logger.error(f"Internal server error: {str(error)}")
    return jsonify({'error': 'Internal server error'}), 500

@app.errorhandler(413)
def too_large(error):
    return jsonify({'error': 'File too large. Maximum size is 50MB.'}), 413

# Cleanup function
def cleanup_old_files():
    """Clean up old uploaded files and plots"""
    try:
        import time
        current_time = time.time()
        max_age = 24 * 60 * 60  # 24 hours
        
        # Clean old uploads
        if os.path.exists(UPLOAD_FOLDER):
            for filename in os.listdir(UPLOAD_FOLDER):
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                if os.path.isfile(file_path):
                    if current_time - os.path.getctime(file_path) > max_age:
                        try:
                            os.unlink(file_path)
                            logger.info(f"🧹 Cleaned old upload: {filename}")
                        except Exception as e:
                            logger.error(f"Error cleaning {file_path}: {str(e)}")
        
        # Clean old plots
        if os.path.exists(PLOTS_FOLDER):
            for filename in os.listdir(PLOTS_FOLDER):
                file_path = os.path.join(PLOTS_FOLDER, filename)
                if os.path.isfile(file_path):
                    if current_time - os.path.getctime(file_path) > max_age:
                        try:
                            os.unlink(file_path)
                            logger.info(f"🧹 Cleaned old plot: {filename}")
                        except Exception as e:
                            logger.error(f"Error cleaning {file_path}: {str(e)}")
                            
    except Exception as e:
        logger.error(f"Error in cleanup: {str(e)}")

# Configure Flask app (runs for both gunicorn and direct execution)
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # 200MB max file size
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
cleanup_old_files()

if __name__ == '__main__':
    logger.info("Backend started")
    logger.info(f"Upload folder: {UPLOAD_FOLDER}")
    logger.info(f"Plots folder: {PLOTS_FOLDER}")
    logger.info(f"Vector DB folder: {VECTOR_DB_DIR}")
    logger.info(f"Gemini API: {'Configured' if GEMINI_API_KEY else 'Not configured'}")
    logger.info(f"Postgres: {'Connected' if using_postgres() else 'Not configured (using local ChromaDB)'}")
    logger.info(f"Server: http://localhost:5000")

    app.run(
        host='0.0.0.0',
        port=5000,
        debug=False,
        threaded=True,
        use_reloader=False
    )